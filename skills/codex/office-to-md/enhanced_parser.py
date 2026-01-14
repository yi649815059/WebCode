"""
Enhanced document parser
Supports better text format recognition and exception handling
Now supports: docx, pdf, xlsx, xls, pptx (all cross-platform formats)
"""

import os
import uuid
import logging
import re
from pathlib import Path
from typing import Dict, List, Tuple, Optional
from datetime import date, datetime, time
import docx
from docx.shared import Inches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import fitz  # PyMuPDF
from PIL import Image
import io

# Import additional parsing libraries
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
    from openpyxl import load_workbook
    from openpyxl.utils.datetime import from_excel
except ImportError:
    openpyxl = None
    from_excel = None

try:
    import xlrd
    from xlrd import xldate_as_datetime
except ImportError:
    xlrd = None
    xldate_as_datetime = None

# Import doc converter
try:
    from doc_converter import DocConverter
except ImportError:
    DocConverter = None

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class EnhancedDocumentParser:
    """Enhanced document parser"""
    
    def __init__(self, image_base_url: str = "http://localhost:5000", 
                 image_save_dir: str = "./static/images",
                 filter_headers_footers: bool = True,
                 enable_doc_conversion: bool = True,
                 libreoffice_path: Optional[str] = None,
                 pdf_table_strategy: str = "lines_strict"):
        """
        Initialize document parser
        
        Args:
            image_base_url: Base URL for image access
            image_save_dir: Directory to save images
            filter_headers_footers: Whether to filter headers and footers (default True)
            enable_doc_conversion: Whether to enable .doc conversion functionality (default True)
            libreoffice_path: LibreOffice executable path (optional)
            pdf_table_strategy: PDF table detection strategy (default "lines_strict")
                - "lines_strict": Strict line detection, fastest, suitable for standard tables
                - "lines": Normal line detection, relatively fast
                - "text": Text position-based detection, slower but more accurate
        """
        self.image_base_url = image_base_url
        self.image_save_dir = Path(image_save_dir)
        self.image_save_dir.mkdir(parents=True, exist_ok=True)
        self.filter_headers_footers = filter_headers_footers
        self.enable_doc_conversion = enable_doc_conversion
        
        # Supported image formats
        self.supported_image_formats = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'}
        
        # Temporary variable to store table block indices
        self._last_table_block_indices = []
        
        # Initialize doc converter
        self.doc_converter = None
        if enable_doc_conversion and DocConverter is not None:
            try:
                self.doc_converter = DocConverter(libreoffice_path)
                logger.info("DOC converter initialized successfully")
            except Exception as e:
                logger.warning(f"DOC converter initialization failed: {str(e)}")
                logger.warning("Will be unable to process .doc format files")
        
    def parse_document(self, file_path: str, cleanup_temp: bool = True) -> Dict:
        """
        Parse document (supports docx, pdf, xlsx, xls, pptx, doc)
        
        Args:
            file_path: Document file path
            cleanup_temp: Whether to clean up temporary files (for .doc conversion)
            
        Returns:
            Dictionary containing markdown content and parsing information
        """
        temp_docx_path = None
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                raise FileNotFoundError(f"File does not exist: {file_path}")
            
            file_ext = file_path.suffix.lower()
            file_size = file_path.stat().st_size
            
            logger.info(f"Starting to parse file: {file_path.name} ({file_size} bytes)")
            
            if file_ext == '.docx':
                return self._parse_docx(str(file_path))
            elif file_ext == '.pdf':
                return self._parse_pdf(str(file_path))
            elif file_ext == '.xlsx':
                return self._parse_xlsx(str(file_path))
            elif file_ext == '.xls':
                return self._parse_xls(str(file_path))
            elif file_ext == '.pptx':
                return self._parse_pptx(str(file_path))
            elif file_ext == '.doc':
                # Use LibreOffice to convert .doc to .docx
                if not self.enable_doc_conversion or self.doc_converter is None:
                    raise ValueError("DOC conversion feature is not enabled or LibreOffice is not installed. Please install LibreOffice or convert the file to .docx format.")
                
                logger.info("Detected .doc file, converting to .docx using LibreOffice...")
                success, result = self.doc_converter.convert_doc_to_docx(str(file_path))
                
                if not success:
                    raise ValueError(f"DOC conversion failed: {result}")
                
                temp_docx_path = result
                logger.info(f"DOC conversion successful, temporary file: {temp_docx_path}")
                
                # Parse the converted docx file
                parse_result = self._parse_docx(temp_docx_path)
                
                # Update file info, mark as original doc file
                parse_result["file_type"] = "doc (converted from docx)"
                parse_result["file_info"]["name"] = file_path.name
                parse_result["file_info"]["original_format"] = "doc"
                parse_result["file_info"]["converted_from"] = temp_docx_path
                
                return parse_result
            elif file_ext == '.ppt':
                # .ppt format no longer supported - only Windows COM component support, not cross-platform
                raise ValueError("The .ppt format is not supported. Please convert the file to .pptx format for cross-platform support.")
            else:
                raise ValueError(f"Unsupported file format: {file_ext}. Supported formats: docx, doc, pdf, xlsx, xls, pptx")
                
        except Exception as e:
            logger.error(f"Document parsing failed: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "markdown": "",
                "images_count": 0,
                "images": [],
                "file_info": {
                    "name": file_path.name if 'file_path' in locals() else "unknown",
                    "size": file_size if 'file_size' in locals() else 0,
                    "format": file_ext if 'file_ext' in locals() else "unknown"
                }
            }
        finally:
            # Clean up temporary files
            if temp_docx_path and cleanup_temp:
                try:
                    if os.path.exists(temp_docx_path):
                        os.remove(temp_docx_path)
                        logger.info(f"Temporary file cleaned up: {temp_docx_path}")
                except Exception as e:
                    logger.warning(f"Failed to clean up temporary file: {str(e)}")
    
    def _parse_docx(self, file_path: str) -> Dict:
        """Parse docx file"""
        try:
            # Open docx file using the correct method
            doc = docx.Document(file_path)
            markdown_content = []
            images_info = []
            
            # Document statistics
            total_paragraphs = len(doc.paragraphs)
            total_tables = len(doc.tables)
            
            logger.info(f"DOCX document contains {total_paragraphs} paragraphs and {total_tables} tables")
            
            # Collect all paragraph text for header/footer detection
            all_paragraphs_text = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            
            # First extract all images, build mapping from rId to image info
            images_info = []
            try:
                images_info = self._extract_docx_images(doc)
                logger.info(f"Extracted {len(images_info)} images from DOCX")
            except Exception as e:
                logger.warning(f"Error extracting DOCX images: {str(e)}")
            
            # Create table index for insertion at original position
            table_index = 0
            processed_tables = set()
            
            # Traverse all elements in the document's original order (paragraphs and tables)
            for element in doc.element.body:
                # Process tables
                if element.tag.endswith('tbl'):
                    # Find corresponding Table object
                    if table_index < len(doc.tables):
                        table = doc.tables[table_index]
                        if table_index not in processed_tables:
                            try:
                                table_md = self._convert_docx_table_to_markdown(table)
                                if table_md:
                                    markdown_content.append(f"\n**Table {table_index + 1}:**\n")
                                    markdown_content.extend(table_md)
                                    markdown_content.append("")
                                processed_tables.add(table_index)
                            except Exception as e:
                                logger.warning(f"Error processing table {table_index+1}: {str(e)}")
                                markdown_content.append(f"\n*[Table {table_index + 1} parsing failed]*\n")
                        table_index += 1
                
                # Process paragraphs
                elif element.tag.endswith('p'):
                    # Find corresponding Paragraph object
                    for paragraph in doc.paragraphs:
                        if paragraph._element == element:
                            try:
                                # Filter headers and footers
                                if self.filter_headers_footers and self._is_header_or_footer(paragraph, all_paragraphs_text):
                                    logger.debug(f"Filtering header/footer content: {paragraph.text[:50] if paragraph.text else ''}")
                                    break
                                
                                # Process paragraph text
                                markdown_lines = self._process_docx_paragraph(paragraph)
                                markdown_content.extend(markdown_lines)
                                
                                # Check if paragraph contains images, if so insert inline
                                paragraph_images = self._get_paragraph_images(paragraph, images_info)
                                if paragraph_images:
                                    for img_info in paragraph_images:
                                        markdown_content.append(f"![Image]({img_info['url']})")
                                        markdown_content.append("")
                                
                            except Exception as e:
                                logger.warning(f"Error processing paragraph: {str(e)}")
                            break
            
            # Clean and format markdown
            result_markdown = self._clean_markdown(markdown_content)
            
            # If no content, add prompt
            if not result_markdown.strip():
                result_markdown = "*Document content is empty or cannot be parsed*"
            
            return {
                "success": True,
                "markdown": result_markdown,
                "images_count": len(images_info),
                "images": images_info,
                "file_type": "docx",
                "error": None,
                "file_info": {
                    "name": Path(file_path).name,
                    "size": Path(file_path).stat().st_size,
                    "paragraphs": total_paragraphs,
                    "tables": total_tables
                }
            }
            
        except Exception as e:
            logger.error(f"DOCX parsing failed: {str(e)}")
            raise
    
    def _is_header_or_footer(self, paragraph, all_paragraphs_text: List[str]) -> bool:
        """检测段落是否为页眉或页脚"""
        text = paragraph.text.strip()
        
        if not text:
            return False
        
        # 1. 检测纯数字(页码)
        if text.isdigit() and len(text) <= 4:
            return True
        
        # 2. 检测常见页码格式
        page_patterns = [
            r'^第\s*\d+\s*页$',  # 第X页
            r'^-\s*\d+\s*-$',     # -X-
            r'^\d+\s*/\s*\d+$',  # X/Y
            r'^Page\s+\d+$',     # Page X
            r'^\d+\s+of\s+\d+$', # X of Y
        ]
        for pattern in page_patterns:
            if re.match(pattern, text, re.IGNORECASE):
                return True
        
        # 3. Detect repeated short text (possibly header)
        # If same text appears more than 3 times in document and is relatively short, likely a header
        if len(text) < 100:
            count = all_paragraphs_text.count(text)
            if count > 3:
                logger.debug(f"Detected repeated text (appears {count} times): {text[:30]}")
                return True
        
        # 4. Detect common header/footer keywords
        footer_keywords = ['\u7248\u6743\u6240\u6709', '\u4fdd\u7559\u6240\u6709\u6743\u5229', 'all rights reserved', 'copyright', 
                          '\u673a\u5bc6', 'confidential', '\u5185\u90e8\u8d44\u6599']
        text_lower = text.lower()
        if any(keyword in text_lower for keyword in footer_keywords) and len(text) < 150:
            return True
        
        return False
    
    def _process_docx_paragraph(self, paragraph) -> List[str]:
        """处理DOCX段落，返回markdown行"""
        lines = []
        text = paragraph.text.strip()
        
        if not text:
            return [""]
        
        # 检查段落样式
        style_name = paragraph.style.name if paragraph.style else ""
        
        # 处理标题
        if style_name.startswith('Heading'):
            level = self._get_heading_level(style_name)
            lines.append(f"{'#' * level} {text}")
            lines.append("")
        elif style_name.lower() in ['title', 'subtitle']:
            lines.append(f"# {text}")
            lines.append("")
        else:
            # 检查段落对齐方式
            alignment = paragraph.alignment
            if alignment == WD_PARAGRAPH_ALIGNMENT.CENTER:
                # 居中文本处理：判断是否像标题
                # 如果文本较短(少于50字符)且不包含句号，可能是标题
                if len(text) < 50 and '。' not in text and '.' not in text:
                    # 作为二级标题处理
                    lines.append(f"## {text}")
                else:
                    # 作为加粗段落处理，表示强调
                    lines.append(f"**{text}**")
            else:
                # 处理列表项（简单检测）
                if re.match(r'^\s*[\d\w]+[.)]\s+', text) or text.strip().startswith(('•', '-', '*')):
                    # 有序或无序列表
                    if re.match(r'^\s*\d+[.)]\s+', text):
                        # 有序列表
                        match = re.match(r'^\s*\d+[.)]\s+(.+)', text)
                        if match:
                            lines.append(f"1. {match.group(1)}")
                    else:
                        # 无序列表
                        cleaned_text = re.sub(r'^\s*[•\-*]\s+', '', text)
                        lines.append(f"- {cleaned_text}")
                else:
                    # 普通段落
                    lines.append(text)
            
            lines.append("")
        
        return lines
    
    def _parse_pdf(self, file_path: str) -> Dict:
        """Parse pdf file"""
        try:
            doc = fitz.open(file_path)
            markdown_content = []
            images_info = []
            
            total_pages = len(doc)
            logger.info(f"PDF document has {total_pages} pages")
            
            for page_num in range(total_pages):
                try:
                    page = doc.load_page(page_num)
                    
                    # Get page text blocks dictionary in one call (cached to avoid repeated calls)
                    page_dict = page.get_text("dict")
                    
                    # Extract tables and get table position information
                    self._last_table_block_indices = []  # Reset
                    tables = self._extract_pdf_tables_optimized(page, page_num, page_dict)
                    
                    # Extract images
                    page_images = self._extract_pdf_images(page, page_num)
                    if page_images:
                        images_info.extend(page_images)
                    
                    # Mix output content in position order (using cached page_dict)
                    if tables:
                        # Get all content blocks and their positions
                        content_blocks = self._get_ordered_content_blocks_optimized(page_dict, tables, page_images)
                        
                        # Output in order
                        for block in content_blocks:
                            if block['type'] == 'text':
                                markdown_content.append(block['content'])
                                markdown_content.append("")
                            elif block['type'] == 'table':
                                markdown_content.extend(block['content'])
                                markdown_content.append("")
                            elif block['type'] == 'image':
                                markdown_content.append(block['content'])
                                markdown_content.append("")
                    else:
                        # When there are no tables, output text and images in normal order
                        text = page.get_text()
                        if text.strip():
                            processed_text = self._process_pdf_text(text)
                            markdown_content.extend(processed_text)
                        
                        # Output images
                        for img_info in page_images:
                            markdown_content.append(f"![Image]({img_info['url']})")
                            markdown_content.append("")
                    
                except Exception as e:
                    logger.warning(f"Error processing PDF page {page_num + 1}: {str(e)}")
                    continue
            
            doc.close()
            
            # Clean and format markdown
            result_markdown = self._clean_markdown(markdown_content)
            
            # If no content, add prompt
            if not result_markdown.strip():
                result_markdown = "*PDF document content is empty or cannot be parsed*"
            
            return {
                "success": True,
                "markdown": result_markdown,
                "images_count": len(images_info),
                "images": images_info,
                "file_type": "pdf",
                "error": None,
                "file_info": {
                    "name": Path(file_path).name,
                    "size": Path(file_path).stat().st_size,
                    "pages": total_pages
                }
            }
            
        except Exception as e:
            logger.error(f"PDF parsing failed: {str(e)}")
            raise
    
    def _process_pdf_text(self, text: str) -> List[str]:
        """处理PDF文本，改进格式识别"""
        lines = []
        
        # 按段落分割
        paragraphs = text.split('\n\n')
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            # 清理多余的换行符，但保留必要的分段
            para = re.sub(r'\n+', ' ', para)
            para = re.sub(r'\s+', ' ', para)
            
            # 检测标题（大写字母开头，相对较短）
            if (len(para) < 100 and 
                (para.isupper() or 
                 re.match(r'^[A-Z][^.!?]*$', para) or
                 re.match(r'^\d+[\.\s]+[A-Z]', para))):
                lines.append(f"## {para}")
                lines.append("")
            else:
                # 检测列表项
                if re.match(r'^\s*[\d\w]+[.)]\s+', para) or para.startswith(('•', '-', '*')):
                    if re.match(r'^\s*\d+[.)]\s+', para):
                        match = re.match(r'^\s*\d+[.)]\s+(.+)', para)
                        if match:
                            lines.append(f"1. {match.group(1)}")
                    else:
                        cleaned_text = re.sub(r'^\s*[•\-*]\s+', '', para)
                        lines.append(f"- {cleaned_text}")
                else:
                    lines.append(para)
                
                lines.append("")
        
        return lines
    
    def _get_paragraph_images(self, paragraph, all_images: List[Dict]) -> List[Dict]:
        """获取段落中包含的图片"""
        paragraph_images = []
        
        try:
            # 检查段落的runs中是否有图片
            for run in paragraph.runs:
                # 检查run中是否包含图片元素
                if 'graphic' in run._element.xml:
                    # 提取图片的关系ID
                    for inline in run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing//{http://schemas.openxmlformats.org/drawingml/2006/picture}blipFill//{http://schemas.openxmlformats.org/drawingml/2006/main}blip'):
                        embed = inline.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if embed:
                            # 查找对应的图片信息
                            for img_info in all_images:
                                if img_info.get('rel_id') == embed:
                                    paragraph_images.append(img_info)
                                    break
        except Exception as e:
            logger.debug(f"获取段落图片时出错: {str(e)}")
        
        return paragraph_images
    
    def _extract_docx_images(self, doc) -> List[Dict]:
        """从docx文档中提取图片"""
        images_info = []
        
        try:
            # 遍历文档中的所有关系，查找图片
            for rel_id, rel in doc.part.rels.items():
                if "image" in rel.target_ref:
                    try:
                        # 获取图片数据
                        image_data = rel.target_part.blob
                        
                        if len(image_data) < 100:  # 跳过太小的图片
                            continue
                        
                        # 生成唯一文件名
                        image_id = str(uuid.uuid4())
                        image_ext = self._get_image_extension(image_data)
                        image_filename = f"{image_id}.{image_ext}"
                        image_path = self.image_save_dir / image_filename
                        
                        # 保存图片
                        with open(image_path, 'wb') as f:
                            f.write(image_data)
                        
                        # 验证图片是否有效
                        if self._validate_image(image_path):
                            # 生成访问URL
                            image_url = f"{self.image_base_url}/static/images/{image_filename}"
                            
                            images_info.append({
                                "filename": image_filename,
                                "path": str(image_path),
                                "url": image_url,
                                "size": len(image_data),
                                "format": image_ext,
                                "rel_id": rel_id  # 保存关系ID用于定位
                            })
                        else:
                            # 删除无效图片
                            os.remove(image_path)
                            
                    except Exception as e:
                        logger.warning(f"处理DOCX图片{rel_id}时出错: {str(e)}")
                        continue
                        
        except Exception as e:
            logger.warning(f"DOCX图片提取失败: {str(e)}")
        
        logger.info(f"从DOCX提取了{len(images_info)}张有效图片")
        return images_info
    
    def _extract_pdf_images(self, page, page_num: int) -> List[Dict]:
        """从pdf页面中提取图片"""
        images_info = []
        
        try:
            image_list = page.get_images()
            logger.debug(f"PDF第{page_num + 1}页发现{len(image_list)}个图片对象")
            
            for img_index, img in enumerate(image_list):
                try:
                    # 获取图片对象
                    xref = img[0]
                    pix = fitz.Pixmap(page.parent, xref)
                    
                    # 跳过CMYK图片和太小的图片
                    if pix.n - pix.alpha >= 4 or pix.width < 50 or pix.height < 50:
                        pix = None
                        continue
                    
                    # 生成唯一文件名
                    image_id = str(uuid.uuid4())
                    image_filename = f"{image_id}.png"
                    image_path = self.image_save_dir / image_filename
                    
                    # 保存图片
                    pix.save(str(image_path))
                    
                    # 验证图片是否有效
                    if self._validate_image(image_path):
                        # 生成访问URL
                        image_url = f"{self.image_base_url}/static/images/{image_filename}"
                        
                        images_info.append({
                            "filename": image_filename,
                            "path": str(image_path),
                            "url": image_url,
                            "page": page_num + 1,
                            "index": img_index,
                            "width": pix.width,
                            "height": pix.height,
                            "format": "png"
                        })
                    else:
                        # 删除无效图片
                        os.remove(image_path)
                    
                    pix = None  # 释放内存
                    
                except Exception as e:
                    logger.warning(f"处理PDF第{page_num + 1}页图片{img_index}时出错: {str(e)}")
                    continue
                    
        except Exception as e:
            logger.warning(f"PDF第{page_num + 1}页图片提取失败: {str(e)}")
        
        logger.debug(f"从PDF第{page_num + 1}页提取了{len(images_info)}张有效图片")
        return images_info
    
    def _extract_pdf_tables_optimized(self, page, page_num: int, page_dict: Dict = None) -> List[Dict]:
        """从PDF页面提取表格，返回表格信息列表（包含内容和位置） - 性能优化版"""
        tables = []
        table_count = 0
        
        try:
            # 方法1: 使用PyMuPDF的find_tables()方法提取标准表格
            # 注意: PyMuPDF的find_tables()不需要strategy参数，已经足够快
            tab_finder = page.find_tables()
            
            if tab_finder and tab_finder.tables:
                logger.info(f"PDF第{page_num + 1}页发现{len(tab_finder.tables)}个标准表格")
                
                for table_idx, table in enumerate(tab_finder.tables):
                    try:
                        table_data = table.extract()
                        if table_data and len(table_data) > 0:
                            markdown_table = self._convert_pdf_table_to_markdown(table_data, page_num, table_count)
                            if markdown_table:
                                # 获取表格的边界框用于定位
                                bbox = table.bbox  # (x0, y0, x1, y1)
                                tables.append({
                                    'content': markdown_table,
                                    'bbox': bbox,
                                    'y_pos': bbox[1],  # 使用上边界的Y坐标
                                    'type': 'standard'  # 标准表格
                                })
                                table_count += 1
                    except Exception as e:
                        logger.warning(f"处理PDF第{page_num + 1}页标准表格{table_idx + 1}时出错: {str(e)}")
                        continue
            
            # 方法2: 如果没找到标准表格,尝试基于文本位置检测表格
            # 优化: 传入已缓存的page_dict，避免重复调用get_text("dict")
            if table_count == 0:
                logger.debug(f"PDF第{page_num + 1}页未找到标准表格,尝试文本位置检测...")
                if page_dict is None:
                    page_dict = page.get_text("dict")
                detected_tables = self._detect_tables_from_text_optimized(page_dict, page_num)
                for detected_table in detected_tables:
                    # detected_table已经是markdown内容列表，需要包装成字典
                    tables.append({
                        'content': detected_table,
                        'bbox': None,
                        'y_pos': 0,  # 文本检测的表格位置已通过_last_table_block_indices处理
                        'type': 'detected'  # 检测到的表格
                    })
                    table_count += 1
                    
        except Exception as e:
            logger.debug(f"PDF第{page_num + 1}页表格提取失败: {str(e)}")
        
        return tables
    
    # 保留旧方法以向后兼容
    def _extract_pdf_tables(self, page, page_num: int) -> List[Dict]:
        """从PDF页面提取表格（兼容旧版本，内部调用优化版）"""
        return self._extract_pdf_tables_optimized(page, page_num, None)
    
    def _detect_tables_from_text_optimized(self, page_dict: Dict, page_num: int) -> List[List[str]]:
        """基于文本位置智能检测表格 - 性能优化版（使用缓存的page_dict）"""
        tables = []
        table_block_indices = []
        
        try:
            # 使用传入的page_dict，避免重复调用
            blocks = page_dict["blocks"]
            
            # 优化: 使用列表推导式和生成器，减少循环层级
            text_blocks = []
            for block_idx, block in enumerate(blocks):
                if block.get("type") == 0:  # 文本块
                    for line in block.get("lines", []):
                        spans = line.get("spans", [])
                        if not spans:
                            continue
                        text = " ".join(span.get("text", "").strip() for span in spans).strip()
                        if text:
                            bbox = line["bbox"]
                            text_blocks.append({
                                "text": text,
                                "x0": bbox[0],
                                "y0": bbox[1],
                                "x1": bbox[2],
                                "y1": bbox[3],
                                "block_idx": block_idx
                            })
            
            if len(text_blocks) < 3:  # 至少需要3行才能构成表格
                return tables
            
            # 按Y坐标排序(从上到下)
            text_blocks.sort(key=lambda b: b["y0"])
            
            # 优化: 使用更高效的行分组算法
            rows = []
            current_row = []
            last_y = None
            y_tolerance = 5
            
            for block in text_blocks:
                if last_y is None or abs(block["y0"] - last_y) < y_tolerance:
                    current_row.append(block)
                    last_y = block["y0"]
                else:
                    if current_row:
                        current_row.sort(key=lambda b: b["x0"])
                        rows.append(current_row)
                    current_row = [block]
                    last_y = block["y0"]
            
            if current_row:
                current_row.sort(key=lambda b: b["x0"])
                rows.append(current_row)
            
            # 优化: 提前退出条件
            if len(rows) < 2:
                return tables
            
            # 统计每行的列数
            col_counts = [len(row) for row in rows]
            if not col_counts:
                return tables
            
            most_common_cols = max(set(col_counts), key=col_counts.count)
            
            # 如果大多数行有相同的列数,认为是表格
            if col_counts.count(most_common_cols) >= len(rows) * 0.6 and most_common_cols > 1:
                # 构建表格数据并记录块索引
                table_data = []
                block_idx_set = set()  # 使用set提升查找性能
                
                for row in rows:
                    if len(row) == most_common_cols:
                        table_data.append([cell["text"] for cell in row])
                        # 记录表格所属的块索引
                        for cell in row:
                            block_idx_set.add(cell["block_idx"])
                
                if table_data:
                    markdown_table = self._convert_pdf_table_to_markdown(table_data, page_num, 0)
                    if markdown_table:
                        tables.append(markdown_table)
                        # 保存表格块索引供后续过滤使用
                        self._last_table_block_indices = list(block_idx_set)
                        logger.info(f"PDF第{page_num + 1}页通过文本位置检测到表格: {len(table_data)}行 x {most_common_cols}列")
                        
        except Exception as e:
            logger.warning(f"PDF第{page_num + 1}页基于文本位置的表格检测失败: {str(e)}")
        
        return tables
    
    # 保留旧方法以向后兼容
    def _detect_tables_from_text(self, page, page_num: int) -> List[List[str]]:
        """基于文本位置智能检测表格（兼容旧版本）"""
        page_dict = page.get_text("dict")
        return self._detect_tables_from_text_optimized(page_dict, page_num)
    
    def _convert_pdf_table_to_markdown(self, table_data: List[List], page_num: int, table_idx: int) -> List[str]:
        """将PDF表格数据转换为markdown格式"""
        if not table_data:
            return []
        
        markdown_table = []
        
        try:
            # 过滤掉完全为空的行
            filtered_data = []
            for row in table_data:
                # 清理单元格数据
                cleaned_row = []
                for cell in row:
                    if cell is None:
                        cleaned_row.append("")
                    else:
                        # 转换为字符串并清理
                        cell_str = str(cell).strip().replace('\n', ' ').replace('|', '\\|')
                        cleaned_row.append(cell_str)
                
                # 如果行中有任何非空内容，则保留
                if any(cell for cell in cleaned_row):
                    filtered_data.append(cleaned_row)
            
            if not filtered_data:
                return []
            
            # 确保所有行有相同的列数
            max_cols = max(len(row) for row in filtered_data)
            normalized_data = []
            for row in filtered_data:
                # 补齐列数
                normalized_row = row + [""] * (max_cols - len(row))
                normalized_data.append(normalized_row)
            
            # 添加表格标题
            markdown_table.append(f"\n**表格 {table_idx + 1}** (第{page_num + 1}页):\n")
            
            # 生成表头
            if normalized_data:
                headers = normalized_data[0]
                # 如果第一行看起来不像表头（全是空或很短的数字），生成默认表头
                if all(not h or (h.isdigit() and len(h) <= 3) for h in headers):
                    headers = [f"列{i+1}" for i in range(max_cols)]
                    # 不跳过第一行，因为它是数据
                    data_rows = normalized_data
                else:
                    # 第一行作为表头
                    data_rows = normalized_data[1:]
                
                # 如果表头为空，使用默认列名
                headers = [h if h else f"列{i+1}" for i, h in enumerate(headers)]
                
                markdown_table.append("| " + " | ".join(headers) + " |")
                markdown_table.append("| " + " | ".join(["---"] * len(headers)) + " |")
                
                # 生成数据行
                for row in data_rows:
                    # 替换空单元格为"-"以提高可读性
                    row_display = [cell if cell else "-" for cell in row]
                    markdown_table.append("| " + " | ".join(row_display) + " |")
            
        except Exception as e:
            logger.warning(f"PDF表格转换失败: {str(e)}")
            return ["*表格解析失败*"]
        
        return markdown_table
    
    def _validate_image(self, image_path: Path) -> bool:
        """验证图片是否有效"""
        try:
            with Image.open(image_path) as img:
                img.verify()
            return True
        except Exception:
            return False
    
    def _get_heading_level(self, style_name: str) -> int:
        """根据样式名称获取标题级别"""
        heading_levels = {
            'Heading 1': 1, 'Title': 1,
            'Heading 2': 2, 'Subtitle': 2,
            'Heading 3': 3,
            'Heading 4': 4,
            'Heading 5': 5,
            'Heading 6': 6
        }
        return heading_levels.get(style_name, 1)
    
    def _convert_docx_table_to_markdown(self, table) -> List[str]:
        """将docx表格转换为markdown格式"""
        markdown_table = []
        
        try:
            rows = table.rows
            if not rows:
                return []
            
            # 检查表格是否有内容
            has_content = False
            for row in rows:
                for cell in row.cells:
                    if cell.text.strip():
                        has_content = True
                        break
                if has_content:
                    break
            
            if not has_content:
                return ["*[空表格]*"]
            
            # 处理表头
            header_row = rows[0]
            headers = [cell.text.strip() or f"列{i+1}" for i, cell in enumerate(header_row.cells)]
            markdown_table.append("| " + " | ".join(headers) + " |")
            markdown_table.append("| " + " | ".join(["---"] * len(headers)) + " |")
            
            # 处理数据行
            for row in rows[1:]:
                row_data = []
                for i, cell in enumerate(row.cells):
                    cell_text = cell.text.strip().replace('\n', '<br>')
                    row_data.append(cell_text or "-")
                markdown_table.append("| " + " | ".join(row_data) + " |")
            
        except Exception as e:
            logger.warning(f"表格转换失败: {str(e)}")
            return ["*[表格解析失败]*"]
        
        return markdown_table
    
    def _get_image_extension(self, image_data: bytes) -> str:
        """根据图片数据判断图片格式"""
        try:
            image = Image.open(io.BytesIO(image_data))
            format_name = image.format.lower()
            if format_name == 'jpeg':
                return 'jpg'
            return format_name.strip() if format_name else 'png'
        except Exception:
            return 'png'  # 默认png格式
    
    def _clean_markdown(self, markdown_lines: List[str]) -> str:
        """清理和格式化markdown内容"""
        # 合并所有行
        content = "\n".join(markdown_lines)
        
        # 清理多余的空行（最多保留2个连续空行）
        content = re.sub(r'\n{3,}', '\n\n', content)
        
        # 清理行首行尾空格
        lines = [line.rstrip() for line in content.split('\n')]
        
        # 移除文档开头和结尾的空行
        while lines and not lines[0].strip():
            lines.pop(0)
        while lines and not lines[-1].strip():
            lines.pop()
        
        return '\n'.join(lines)
    
    def _get_ordered_content_blocks_optimized(self, page_dict: Dict, tables, images) -> List[Dict]:
        """获取按位置排序的内容块(文本、表格、图片) - 性能优化版"""
        content_blocks = []
        
        try:
            # 使用传入的page_dict，避免重复调用
            blocks = page_dict["blocks"]
            image_idx = 0
            
            # 优化: 预先筛选标准表格，避免重复过滤
            standard_tables = [t for t in tables if t.get('type') == 'standard' and t.get('bbox')]
            detected_tables = [t for t in tables if t.get('type') == 'detected']
            
            # 添加标准表格
            for table_info in standard_tables:
                y_pos = table_info['y_pos']
                content_blocks.append({
                    'type': 'table',
                    'content': table_info['content'],
                    'y_pos': y_pos,
                    'sort_key': (y_pos, -1)
                })
            
            # 优化: 使用集合快速查找
            table_block_indices_set = set(self._last_table_block_indices)
            min_table_block_idx = min(self._last_table_block_indices) if self._last_table_block_indices else -1
            
            # 处理文本块和检测到的表格
            for block_idx, block in enumerate(blocks):
                block_type = block.get("type")
                
                if block_type == 0:  # 文本块
                    bbox = block.get("bbox", [0, 0, 0, 0])
                    y_pos = bbox[1]
                    
                    # 优化: 只在有标准表格时才检查重叠
                    is_in_table = False
                    if standard_tables:
                        for table_info in standard_tables:
                            t_bbox = table_info['bbox']
                            if (bbox[0] >= t_bbox[0] - 5 and bbox[2] <= t_bbox[2] + 5 and
                                bbox[1] >= t_bbox[1] - 5 and bbox[3] <= t_bbox[3] + 5):
                                is_in_table = True
                                break
                    
                    if is_in_table:
                        continue
                    
                    # 优化: 使用集合查找代替列表查找
                    if block_idx in table_block_indices_set:
                        # 只在第一个表格块处插入检测到的表格
                        if block_idx == min_table_block_idx and detected_tables:
                            content_blocks.append({
                                'type': 'table',
                                'content': detected_tables[0]['content'],
                                'y_pos': y_pos,
                                'sort_key': (y_pos, block_idx)
                            })
                    else:
                        # 优化: 简化文本提取逻辑
                        text_parts = []
                        for line in block.get("lines", []):
                            spans = line.get("spans", [])
                            if spans:
                                text = " ".join(span.get("text", "").strip() for span in spans).strip()
                                if text:
                                    text_parts.append(text)
                        
                        if text_parts:
                            content_blocks.append({
                                'type': 'text',
                                'content': " ".join(text_parts),
                                'y_pos': y_pos,
                                'sort_key': (y_pos, block_idx)
                            })
                
                elif block_type == 1:  # 图片块
                    if image_idx < len(images):
                        bbox = block.get("bbox", [0, 0, 0, 0])
                        y_pos = bbox[1]
                        img_info = images[image_idx]
                        content_blocks.append({
                            'type': 'image',
                            'content': f"![图片]({img_info['url']})",
                            'y_pos': y_pos,
                            'sort_key': (y_pos, block_idx)
                        })
                        image_idx += 1
            
            # 按Y坐标和块索引排序
            content_blocks.sort(key=lambda x: x['sort_key'])
            
        except Exception as e:
            logger.warning(f"获取有序内容块失败: {str(e)}")
        
        return content_blocks
    
    # 保留旧方法以向后兼容
    def _get_ordered_content_blocks(self, page, tables, images) -> List[Dict]:
        """获取按位置排序的内容块（兼容旧版本）"""
        page_dict = page.get_text("dict")
        return self._get_ordered_content_blocks_optimized(page_dict, tables, images)
    
    def _extract_non_table_text(self, page) -> str:
        """提取非表格区域的文本"""
        try:
            blocks = page.get_text("dict")["blocks"]
            non_table_texts = []
            
            for block_idx, block in enumerate(blocks):
                # 跳过表格块
                if block_idx in self._last_table_block_indices:
                    continue
                
                if block.get("type") == 0:  # 文本块
                    for line in block.get("lines", []):
                        spans = line.get("spans", [])
                        if spans:
                            text = " ".join(span.get("text", "").strip() for span in spans).strip()
                            if text:
                                non_table_texts.append(text)
            
            return " ".join(non_table_texts) if non_table_texts else ""
        except Exception as e:
            logger.warning(f"提取非表格文本失败: {str(e)}")
            return ""
    
    def _get_table_data_from_markdown(self, table_markdown: List[str]) -> int:
        """从markdown表格中获取数据行数"""
        # 过滤掉表格标题、分隔线等,只计算数据行
        data_rows = 0
        for line in table_markdown:
            if line.strip().startswith('|') and '---' not in line and '**表格' not in line:
                data_rows += 1
        return max(0, data_rows - 1)  # 减去表头行

    def _format_dt(self, dt: date | datetime) -> str:
        """将日期/时间统一转为可读文本。"""
        if isinstance(dt, datetime):
            if dt.time() == time(0, 0):
                return dt.date().isoformat()
            return dt.strftime("%Y-%m-%d %H:%M:%S")
        if isinstance(dt, date):
            return dt.isoformat()
        return str(dt)

    def _cell_to_text(self, cell, wb_epoch) -> str:
        """尽量按单元格格式输出，修复日期显示为数字的问题。"""
        v = cell.value
        if v is None:
            return ""

        # 直接标记为日期的单元格
        if cell.is_date:
            if isinstance(v, (int, float)):
                try:
                    v = from_excel(v, wb_epoch)
                except Exception:
                    pass
            return self._format_dt(v)

        # 数值但格式看起来像日期/时间
        if isinstance(v, (int, float)):
            fmt = (cell.number_format or "").lower()
            if any(token in fmt for token in ("yy", "dd", "mm", "hh", "ss")):
                try:
                    dt = from_excel(v, wb_epoch)
                    return self._format_dt(dt)
                except Exception:
                    pass

            # 兜底：部分表格日期存为整数但格式为常规，尝试按序列号转换
            if float(v).is_integer() and 20000 <= v <= 60000:  # 约对应 1955-2050 之间
                try:
                    dt = from_excel(v, wb_epoch)
                    if isinstance(dt, (date, datetime)) and date(1990, 1, 1) <= getattr(dt, "date", lambda: dt)() <= date(2100, 12, 31):
                        return self._format_dt(dt)
                except Exception:
                    pass

        return str(v)

    def _parse_xlsx(self, file_path: str) -> Dict:
        """解析xlsx文件"""
        try:
            if openpyxl is None:
                raise ImportError("需要安装openpyxl库来解析xlsx文件")
            
            workbook = load_workbook(filename=file_path, data_only=True)
            wb_epoch = workbook.epoch
            markdown_content = []
            images_info = []
            
            logger.info(f"XLSX文件包含{len(workbook.sheetnames)}个工作表")
            
            # 提取图片
            try:
                images_info = self._extract_xlsx_images(file_path)
                if images_info:
                    logger.info(f"从XLSX文件中提取了{len(images_info)}张图片")
            except Exception as e:
                logger.warning(f"XLSX图片提取失败: {str(e)}")
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                # 移除工作表标题，直接处理数据
                
                # 获取有数据的范围
                if sheet.max_row == 1 and sheet.max_column == 1:
                    # 空工作表跳过，不添加提示
                    continue
                
                # 转换为表格格式
                table_data = []
                for row in sheet.iter_rows(min_row=1, max_row=sheet.max_row, 
                                         min_col=1, max_col=sheet.max_column):
                    # 使用新的单元格处理函数，支持日期格式
                    row_values = [self._cell_to_text(cell, wb_epoch) for cell in row]
                    
                    # 过滤掉完全为空的行
                    if any(cell.strip() for cell in row_values):
                        table_data.append(row_values)
                
                if table_data:
                    markdown_content.extend(self._convert_list_to_markdown_table(table_data))
                    markdown_content.append("")
                    
                    # 在每个工作表后展示对应的图片(如果有)
                    sheet_images = [img for img in images_info if img.get('sheet') == sheet_name]
                    if sheet_images:
                        for img_info in sheet_images:
                            markdown_content.append(f"![{img_info['filename']}]({img_info['url']})")
                        markdown_content.append("")
            
            result_markdown = self._clean_markdown(markdown_content)
            
            return {
                "success": True,
                "markdown": result_markdown,
                "images_count": len(images_info),
                "images": images_info,
                "file_type": "xlsx",
                "error": None,
                "file_info": {
                    "name": Path(file_path).name,
                    "size": Path(file_path).stat().st_size,
                    "sheets": len(workbook.sheetnames)
                }
            }
            
        except Exception as e:
            logger.error(f"XLSX解析失败: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "markdown": "",
                "images_count": 0,
                "images": [],
                "file_type": "xlsx"
            }

    def _xls_cell_to_text(self, cell, datemode) -> str:
        """处理xls单元格，支持日期格式转换。"""
        if cell.value is None or cell.value == "":
            return ""
        
        # xlrd中，ctype=3表示日期类型
        if cell.ctype == 3:  # XL_CELL_DATE
            try:
                dt = xldate_as_datetime(cell.value, datemode)
                return self._format_dt(dt)
            except Exception:
                pass
        
        return str(cell.value) if cell.value is not None else ""

    def _parse_xls(self, file_path: str) -> Dict:
        """解析xls文件"""
        try:
            if xlrd is None:
                raise ImportError("需要安装xlrd库来解析xls文件")
            
            workbook = xlrd.open_workbook(file_path, formatting_info=False)
            datemode = workbook.datemode
            markdown_content = []
            images_info = []
            
            logger.info(f"XLS文件包含{workbook.nsheets}个工作表")
            
            # 提取图片 - 注意：xls格式的图片提取比较困难，主要支持xlsx
            try:
                images_info = self._extract_xls_images(file_path)
                if images_info:
                    logger.info(f"从XLS文件中提取了{len(images_info)}张图片")
            except Exception as e:
                logger.warning(f"XLS图片提取失败: {str(e)}")
            
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                sheet_name = workbook.sheet_names()[sheet_idx]
                # 移除工作表标题，直接处理数据
                
                if sheet.nrows == 0:
                    # 空工作表跳过，不添加提示
                    continue
                
                # 转换为表格格式
                table_data = []
                for row_idx in range(sheet.nrows):
                    row_data = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        # 使用新的单元格处理函数，支持日期格式
                        cell_value = self._xls_cell_to_text(cell, datemode)
                        row_data.append(cell_value)
                    
                    # 过滤掉完全为空的行
                    if any(cell.strip() for cell in row_data):
                        table_data.append(row_data)
                
                if table_data:
                    markdown_content.extend(self._convert_list_to_markdown_table(table_data))
                    markdown_content.append("")
                    
                    # 在每个工作表后展示对应的图片(如果有)
                    sheet_images = [img for img in images_info if img.get('sheet') == sheet_name]
                    if sheet_images:
                        for img_info in sheet_images:
                            markdown_content.append(f"![{img_info['filename']}]({img_info['url']})")
                        markdown_content.append("")
            
            result_markdown = self._clean_markdown(markdown_content)
            
            return {
                "success": True,
                "markdown": result_markdown,
                "images_count": len(images_info),
                "images": images_info,
                "file_type": "xls",
                "error": None,
                "file_info": {
                    "name": Path(file_path).name,
                    "size": Path(file_path).stat().st_size,
                    "sheets": workbook.nsheets
                }
            }
            
        except Exception as e:
            logger.error(f"XLS解析失败: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "markdown": "",
                "images_count": 0,
                "images": [],
                "file_type": "xls"
            }

    def _parse_pptx(self, file_path: str) -> Dict:
        """解析pptx文件"""
        try:
            if Presentation is None:
                raise ImportError("需要安装python-pptx库来解析pptx文件")
            
            prs = Presentation(file_path)
            markdown_content = []
            images_info = []
            
            logger.info(f"PPTX文件包含{len(prs.slides)}张幻灯片")
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                # 移除幻灯片标题，直接提取内容
                
                # 提取文本内容
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # 提取图片
                    if shape.shape_type == 13:  # PICTURE = 13
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            image_ext = self._get_image_extension(image_bytes)
                            image_filename = f"{uuid.uuid4()}.{image_ext}"
                            image_path = self.image_save_dir / image_filename
                            
                            with open(image_path, 'wb') as f:
                                f.write(image_bytes)
                            
                            image_url = f"{self.image_base_url}/static/images/{image_filename}"
                            images_info.append({
                                "filename": image_filename,
                                "url": image_url,
                                "size": len(image_bytes)
                            })
                            
                            slide_text.append(f"![图片]({image_url})")
                            
                        except Exception as e:
                            logger.warning(f"提取PPT图片失败: {str(e)}")
                
                if slide_text:
                    markdown_content.extend(slide_text)
                else:
                    markdown_content.append("*幻灯片无文本内容*")
                
                markdown_content.append("")
            
            result_markdown = self._clean_markdown(markdown_content)
            
            return {
                "success": True,
                "markdown": result_markdown,
                "images_count": len(images_info),
                "images": images_info,
                "file_type": "pptx",
                "error": None,
                "file_info": {
                    "name": Path(file_path).name,
                    "size": Path(file_path).stat().st_size,
                    "slides": len(prs.slides)
                }
            }
            
        except Exception as e:
            logger.error(f"PPTX解析失败: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "markdown": "",
                "images_count": 0,
                "images": [],
                "file_type": "pptx"
            }

    def _convert_list_to_markdown_table(self, table_data: List[List[str]]) -> List[str]:
        """将二维列表转换为markdown表格"""
        if not table_data:
            return ["*空表格*"]
        
        markdown_table = []
        
        try:
            # 确保所有行有相同的列数
            max_cols = max(len(row) for row in table_data)
            normalized_data = []
            for row in table_data:
                # 补齐列数
                normalized_row = row + [""] * (max_cols - len(row))
                # 清理单元格内容
                cleaned_row = [str(cell).replace('\n', '<br>').replace('|', '\\|') for cell in normalized_row]
                normalized_data.append(cleaned_row)
            
            # 生成表头
            if normalized_data:
                headers = normalized_data[0] if len(normalized_data) > 1 else [f"列{i+1}" for i in range(max_cols)]
                markdown_table.append("| " + " | ".join(headers) + " |")
                markdown_table.append("| " + " | ".join(["---"] * len(headers)) + " |")
                
                # 生成数据行
                data_rows = normalized_data[1:] if len(normalized_data) > 1 else normalized_data
                for row in data_rows:
                    markdown_table.append("| " + " | ".join(row) + " |")
            
        except Exception as e:
            logger.warning(f"表格转换失败: {str(e)}")
            return ["*表格解析失败*"]
        
        return markdown_table

    def _extract_xlsx_images(self, file_path: str) -> List[Dict]:
        """从xlsx文件中提取图片"""
        images_info = []
        
        try:
            import zipfile
            from xml.etree import ElementTree as ET
            
            # xlsx文件本质上是一个zip文件
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # 查找所有媒体文件
                media_files = [name for name in zip_file.namelist() 
                              if name.startswith('xl/media/') and 
                              any(name.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'])]
                
                logger.info(f"在XLSX文件中发现{len(media_files)}个媒体文件")
                
                # 尝试解析工作表和图片的关系(可选,如果失败就不添加sheet信息)
                sheet_image_mapping = {}
                try:
                    # 读取工作表关系文件
                    for sheet_rel in [name for name in zip_file.namelist() if 'xl/worksheets/_rels/' in name]:
                        sheet_name = sheet_rel.split('/')[-1].replace('.rels', '')
                        content = zip_file.read(sheet_rel)
                        # 这里可以解析XML获取图片关系,但较复杂,暂时简化处理
                except:
                    pass
                
                for media_file in media_files:
                    try:
                        # 读取图片数据
                        image_data = zip_file.read(media_file)
                        
                        if len(image_data) < 100:  # 跳过太小的图片
                            continue
                        
                        # 生成唯一文件名
                        image_id = str(uuid.uuid4())
                        original_ext = Path(media_file).suffix.lower().lstrip('.')
                        original_ext = original_ext.strip()
                        image_filename = f"{image_id}.{original_ext}"
                        image_path = self.image_save_dir / image_filename
                        
                        # 保存图片
                        with open(image_path, 'wb') as f:
                            f.write(image_data)
                        
                        # 验证图片是否有效
                        if self._validate_image(image_path):
                            # 生成访问URL
                            image_url = f"{self.image_base_url}/static/images/{image_filename}"
                            
                            images_info.append({
                                "filename": image_filename,
                                "path": str(image_path),
                                "url": image_url,
                                "size": len(image_data),
                                "format": original_ext,
                                "source": media_file
                            })
                            logger.info(f"成功提取图片: {media_file}")
                        else:
                            # 删除无效图片
                            os.remove(image_path)
                            logger.warning(f"无效图片已删除: {media_file}")
                            
                    except Exception as e:
                        logger.warning(f"处理XLSX图片{media_file}时出错: {str(e)}")
                        
        except Exception as e:
            logger.error(f"XLSX图片提取失败: {str(e)}")
            
        return images_info

    def _extract_xls_images(self, file_path: str) -> List[Dict]:
        """从xls文件中提取图片 - 注意：xls格式图片提取相对困难"""
        images_info = []
        
        try:
            # 对于xls文件，图片提取比较复杂，因为它不是基于XML的格式
            # 这里我们尝试一个简单的方法，但成功率可能不高
            logger.warning("XLS格式的图片提取功能有限，建议使用XLSX格式以获得更好的图片支持")
            
            # 尝试使用openpyxl转换（如果可能的话）
            # 但通常xls文件需要先转换为xlsx才能很好地提取图片
            
        except Exception as e:
            logger.error(f"XLS图片提取失败: {str(e)}")
            
        return images_info


# 向后兼容的别名
DocumentParser = EnhancedDocumentParser
